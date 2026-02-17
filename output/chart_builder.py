"""
Chart and table builder functions for PPTX slides.

Each function adds a visual element (chart or table) to an existing slide.
Functions accept a pptx Slide, a pandas DataFrame, and position/size params.

For native charts: uses python-pptx ChartData objects.
For heatmaps: uses Table shapes with manual conditional cell coloring,
since python-pptx has no native heatmap support.
"""

import logging
from typing import Optional

import pandas as pd
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
from pptx.util import Inches, Pt, Emu

from output.style import (
    COLORS_CATEGORICAL,
    COLORS_BINARY,
    COLOR_BODY_TEXT,
    COLOR_WHITE,
    COLOR_LIGHT_GRAY,
    FONT_BODY,
    SIZE_CHART_LABEL,
    SIZE_TABLE_BODY,
    SIZE_TABLE_HEADER,
    apply_chart_style,
    interpolate_heatmap_color,
    style_table_header,
    style_table_body,
)

logger = logging.getLogger(__name__)


def add_grouped_bar_chart(
    slide,
    data: pd.DataFrame,
    groupby_col: str,
    category_col: str,
    value_col: str,
    left: Inches,
    top: Inches,
    width: Inches,
    height: Inches,
    title: str = "",
    colors: list[RGBColor] | None = None,
) -> None:
    """
    Add a clustered column chart to a slide.

    Pivots the DataFrame so each unique value of category_col becomes a
    separate series, with groupby_col values as the x-axis categories.

    Used by: Slide 1 (Market Fingerprint, x5), Slide 3 (chart portion).

    Args:
        slide: pptx Slide object
        data: DataFrame with at least [groupby_col, category_col, value_col]
        groupby_col: Column for x-axis categories (e.g., "Retailer")
        category_col: Column whose unique values become chart series
        value_col: Column with numeric values to plot
        left, top, width, height: Chart position and size
        title: Optional chart title displayed above the chart
        colors: Series colors. Defaults to COLORS_CATEGORICAL.
    """
    if data.empty:
        _add_no_data_placeholder(slide, left, top, width, height, title)
        return

    if colors is None:
        colors = COLORS_CATEGORICAL

    chart_data = CategoryChartData()

    # Categories = unique values of groupby_col (e.g., each retailer)
    categories = data[groupby_col].unique().tolist()
    chart_data.categories = categories

    # Each unique value in category_col becomes a series
    series_values = data[category_col].unique().tolist()

    for series_name in series_values:
        series_mask = data[category_col] == series_name
        series_df = data[series_mask].set_index(groupby_col)[value_col]

        # Align values to the category order, filling missing with 0
        values = [
            float(series_df.get(cat, 0.0))
            for cat in categories
        ]
        chart_data.add_series(str(series_name), values)

    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        left, top, width, height,
        chart_data,
    )
    chart = chart_frame.chart

    # Set chart title
    if title:
        chart.has_title = True
        chart.chart_title.text_frame.paragraphs[0].text = title
        chart.chart_title.text_frame.paragraphs[0].font.size = SIZE_CHART_LABEL
        chart.chart_title.text_frame.paragraphs[0].font.name = FONT_BODY
        chart.chart_title.text_frame.paragraphs[0].font.bold = True

    # Legend below chart to save vertical space
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    chart.legend.font.size = Pt(8)

    apply_chart_style(chart, colors)


def add_bar_chart(
    slide,
    data: pd.DataFrame,
    category_col: str,
    value_col: str,
    left: Inches,
    top: Inches,
    width: Inches,
    height: Inches,
    title: str = "",
    colors: list[RGBColor] | None = None,
    show_data_labels: bool = True,
    number_format: str = "0.0",
) -> None:
    """
    Add a single-series bar chart to a slide.

    Each row in the DataFrame becomes a bar. Bars can be individually
    colored by passing a list of colors matching the number of categories.

    Used by: Slides 4-10 (retailer deep dives, 4 charts each).

    Args:
        slide: pptx Slide object
        data: DataFrame with [category_col, value_col]
        category_col: Column for bar labels (x-axis)
        value_col: Column with numeric values
        left, top, width, height: Chart position and size
        title: Optional chart title
        colors: Per-bar colors. Defaults to COLORS_CATEGORICAL.
        show_data_labels: Whether to show values above bars
        number_format: Number format for data labels (e.g., "0.0" or "0.0%")
    """
    if data.empty:
        _add_no_data_placeholder(slide, left, top, width, height, title)
        return

    if colors is None:
        colors = COLORS_CATEGORICAL

    chart_data = CategoryChartData()
    categories = data[category_col].tolist()
    values = data[value_col].tolist()

    chart_data.categories = [str(c) for c in categories]
    chart_data.add_series("Value", [float(v) for v in values])

    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        left, top, width, height,
        chart_data,
    )
    chart = chart_frame.chart

    # Title
    if title:
        chart.has_title = True
        chart.chart_title.text_frame.paragraphs[0].text = title
        chart.chart_title.text_frame.paragraphs[0].font.size = SIZE_CHART_LABEL
        chart.chart_title.text_frame.paragraphs[0].font.name = FONT_BODY
        chart.chart_title.text_frame.paragraphs[0].font.bold = True

    # No legend needed for single series
    chart.has_legend = False

    # Color each bar individually
    series = chart.series[0]
    for point_idx in range(len(categories)):
        point = series.points[point_idx]
        point.format.fill.solid()
        point.format.fill.fore_color.rgb = colors[point_idx % len(colors)]

    # Data labels
    if show_data_labels:
        series.has_data_labels = True
        data_labels = series.data_labels
        data_labels.font.size = SIZE_CHART_LABEL
        data_labels.font.name = FONT_BODY
        data_labels.number_format = number_format
        data_labels.number_format_is_linked = False
        data_labels.label_position = XL_LABEL_POSITION.OUTSIDE_END

    # Style axes
    chart.value_axis.tick_labels.font.size = SIZE_CHART_LABEL
    chart.value_axis.tick_labels.font.name = FONT_BODY
    chart.category_axis.tick_labels.font.size = Pt(9)
    chart.category_axis.tick_labels.font.name = FONT_BODY


def add_heatmap_table(
    slide,
    data: pd.DataFrame,
    left: Inches,
    top: Inches,
    width: Inches,
    height: Inches,
    heatmap_columns: list[str],
    format_as_pct: bool = True,
) -> None:
    """
    Add a styled table with conditional heatmap coloring on selected columns.

    Non-heatmap columns are rendered without conditional fill. Header row
    is styled with brand colors. Used for Slide 2 (Brand Landscape).

    Args:
        slide: pptx Slide object
        data: DataFrame to render as a table
        left, top, width, height: Table position and size
        heatmap_columns: Column names that should get heatmap coloring
        format_as_pct: If True, format heatmap values as "X.X%"
    """
    if data.empty:
        _add_no_data_placeholder(slide, left, top, width, height, "Brand Landscape")
        return

    num_rows = len(data) + 1  # +1 for header
    num_cols = len(data.columns)

    table_shape = slide.shapes.add_table(num_rows, num_cols, left, top, width, height)
    table = table_shape.table

    # Distribute column widths proportionally
    # Brand column gets more space, numeric columns get less
    col_width = int(width / num_cols)
    for col_idx in range(num_cols):
        col_name = data.columns[col_idx]
        if col_name == "Brand":
            table.columns[col_idx].width = int(col_width * 1.5)
        else:
            table.columns[col_idx].width = col_width

    # Populate header row
    for col_idx, col_name in enumerate(data.columns):
        cell = table.cell(0, col_idx)
        cell.text = str(col_name)

    # Populate data rows
    for row_idx, (_, row) in enumerate(data.iterrows(), start=1):
        for col_idx, col_name in enumerate(data.columns):
            cell = table.cell(row_idx, col_idx)
            raw_value = row[col_name]

            # Format the display value
            if col_name in heatmap_columns and pd.notna(raw_value):
                display_value = f"{float(raw_value):.1f}%" if format_as_pct else f"{raw_value}"
            elif col_name == "Total Market Share" and pd.notna(raw_value):
                display_value = f"{float(raw_value):.1f}%"
            elif col_name in ("% Cold Pressed", "% Functional") and pd.notna(raw_value):
                display_value = f"{float(raw_value):.1f}%"
            elif pd.notna(raw_value):
                display_value = str(raw_value)
            else:
                display_value = "—"

            cell.text = display_value

    # Apply heatmap coloring to the relevant columns
    for col_name in heatmap_columns:
        if col_name not in data.columns:
            continue

        col_idx = list(data.columns).index(col_name)
        col_values = pd.to_numeric(data[col_name], errors="coerce").fillna(0.0)
        min_val = col_values.min()
        max_val = col_values.max()

        for row_idx in range(len(data)):
            cell = table.cell(row_idx + 1, col_idx)  # +1 to skip header
            cell_value = float(col_values.iloc[row_idx])
            bg_color = interpolate_heatmap_color(cell_value, min_val, max_val)
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg_color

    # Style header and body
    style_table_header(table)
    style_table_body(table)


def add_data_table(
    slide,
    data: pd.DataFrame,
    left: Inches,
    top: Inches,
    width: Inches,
    height: Inches,
    column_widths: list[int] | None = None,
) -> None:
    """
    Add a simple styled table to a slide (no conditional coloring).

    Used by: Slide 3 (Retailer Sizing table portion).

    Args:
        slide: pptx Slide object
        data: DataFrame to render
        left, top, width, height: Table position and size
        column_widths: Optional list of widths (in EMU) per column
    """
    if data.empty:
        _add_no_data_placeholder(slide, left, top, width, height, "Retailer Sizing")
        return

    num_rows = len(data) + 1  # +1 for header
    num_cols = len(data.columns)

    table_shape = slide.shapes.add_table(num_rows, num_cols, left, top, width, height)
    table = table_shape.table

    # Set column widths
    if column_widths:
        for col_idx, col_width in enumerate(column_widths):
            if col_idx < num_cols:
                table.columns[col_idx].width = col_width
    else:
        default_width = int(width / num_cols)
        for col_idx in range(num_cols):
            table.columns[col_idx].width = default_width

    # Populate header
    for col_idx, col_name in enumerate(data.columns):
        table.cell(0, col_idx).text = str(col_name)

    # Populate data
    for row_idx, (_, row) in enumerate(data.iterrows(), start=1):
        for col_idx, col_name in enumerate(data.columns):
            raw_value = row[col_name]
            display_value = str(raw_value) if pd.notna(raw_value) else "—"
            table.cell(row_idx, col_idx).text = display_value

    # Style
    style_table_header(table)
    style_table_body(table)

    # Alternate row shading for readability
    for row_idx in range(2, num_rows, 2):
        for cell in table.rows[row_idx].cells:
            cell.fill.solid()
            cell.fill.fore_color.rgb = COLOR_LIGHT_GRAY


def _add_no_data_placeholder(
    slide,
    left: Inches,
    top: Inches,
    width: Inches,
    height: Inches,
    title: str = "",
) -> None:
    """
    Add a placeholder text box when no data is available for a chart area.

    Args:
        slide: pptx Slide object
        left, top, width, height: Position where the chart would have been
        title: Name of the missing chart for context
    """
    from pptx.enum.text import PP_ALIGN

    text_box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = text_box.text_frame
    text_frame.word_wrap = True

    paragraph = text_frame.paragraphs[0]
    label = f"Insufficient data for: {title}" if title else "Insufficient data"
    paragraph.text = label
    paragraph.font.name = FONT_BODY
    paragraph.font.size = Pt(14)
    paragraph.font.color.rgb = COLOR_BODY_TEXT
    paragraph.font.italic = True
    paragraph.alignment = PP_ALIGN.CENTER

    logger.warning(f"No data available for chart: {title}")
