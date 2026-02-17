"""
PPTX style constants and reusable formatting helpers.

Single source of truth for all visual styling in the PowerPoint output:
fonts, colors, sizes, positions, and helper functions that apply them
consistently across every slide.

Colors are sourced from the Fruity Line template spec and
config/storyline.py CHART_COLORS.
"""

import logging
from pathlib import Path

from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

logger = logging.getLogger(__name__)


# ---------------------------------------------------------------------------
# Slide dimensions (widescreen 16:9)
# ---------------------------------------------------------------------------
SLIDE_WIDTH = Inches(13.33)
SLIDE_HEIGHT = Inches(7.5)

# ---------------------------------------------------------------------------
# Default asset paths
# ---------------------------------------------------------------------------
DEFAULT_TEMPLATE_PATH = Path("assets/FL template.pptx")
DEFAULT_LOGO_PATH = Path("assets/fruity_line_logo.png")

# ---------------------------------------------------------------------------
# Fonts
# ---------------------------------------------------------------------------
FONT_TITLE = "Aptos"
FONT_BODY = "Aptos"
FONT_SOURCE = "Aptos"

# ---------------------------------------------------------------------------
# Font sizes
# ---------------------------------------------------------------------------
SIZE_TITLE = Pt(28)
SIZE_SUBTITLE = Pt(14)
SIZE_BODY = Pt(12)
SIZE_CHART_LABEL = Pt(10)
SIZE_TABLE_HEADER = Pt(11)
SIZE_TABLE_BODY = Pt(10)
SIZE_SOURCE = Pt(8)
SIZE_SLIDE_NUMBER = Pt(8)

# ---------------------------------------------------------------------------
# Colors — template spec
# ---------------------------------------------------------------------------
COLOR_TITLE = RGBColor(0xC0, 0x00, 0x00)       # dark red  #C00000
COLOR_POSITIVE = RGBColor(0x68, 0xB4, 0x45)     # accent green
COLOR_NEGATIVE = RGBColor(0xC0, 0x00, 0x00)     # accent red
COLOR_BODY_TEXT = RGBColor(0x2C, 0x3E, 0x50)     # dark gray
COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_LIGHT_GRAY = RGBColor(0xEC, 0xF0, 0xF1)

# ---------------------------------------------------------------------------
# Colors — brand palette (from storyline.py BRAND_COLORS)
# ---------------------------------------------------------------------------
COLOR_BRAND_PRIMARY = RGBColor(0xFF, 0x6B, 0x35)   # Fruity Line orange
COLOR_BRAND_SECONDARY = RGBColor(0x00, 0x4E, 0x89) # Deep blue
COLOR_BRAND_ACCENT = RGBColor(0xF7, 0xB8, 0x01)    # Golden yellow

# ---------------------------------------------------------------------------
# Chart color palettes (pre-converted to RGBColor)
# ---------------------------------------------------------------------------
COLORS_CATEGORICAL: list[RGBColor] = [
    RGBColor(0xFF, 0x6B, 0x35),  # Orange
    RGBColor(0x00, 0x4E, 0x89),  # Blue
    RGBColor(0xF7, 0xB8, 0x01),  # Yellow
    RGBColor(0x1B, 0x99, 0x8B),  # Teal
    RGBColor(0xC7, 0x3E, 0x1D),  # Red
    RGBColor(0x6A, 0x4C, 0x93),  # Purple
]

COLORS_BINARY: list[RGBColor] = [
    RGBColor(0xFF, 0x6B, 0x35),  # Orange
    RGBColor(0x00, 0x4E, 0x89),  # Blue
]

# Heatmap gradient stops: white → light peach → medium peach → orange → red
HEATMAP_STOPS: list[tuple[float, tuple[int, int, int]]] = [
    (0.0,  (255, 255, 255)),    # White
    (0.25, (255, 229, 217)),    # Light peach
    (0.50, (255, 202, 176)),    # Medium peach
    (0.75, (255, 107, 53)),     # Orange
    (1.0,  (199, 62, 29)),      # Red
]

# ---------------------------------------------------------------------------
# Layout positions (reusable across slides)
# ---------------------------------------------------------------------------
TITLE_LEFT = Inches(0.5)
TITLE_TOP = Inches(0.3)
TITLE_WIDTH = Inches(10.0)
TITLE_HEIGHT = Inches(0.6)

SOURCE_LEFT = Inches(0.5)
SOURCE_TOP = Inches(7.0)
SOURCE_WIDTH = Inches(8.0)
SOURCE_HEIGHT = Inches(0.35)

LOGO_WIDTH = Inches(1.2)
LOGO_HEIGHT = Inches(0.6)
LOGO_LEFT = Inches(11.8)
LOGO_TOP = Inches(6.75)

SLIDE_NUMBER_LEFT = Inches(12.5)
SLIDE_NUMBER_TOP = Inches(7.05)
SLIDE_NUMBER_WIDTH = Inches(0.6)
SLIDE_NUMBER_HEIGHT = Inches(0.3)

# Body content area (below title, above source)
CONTENT_TOP = Inches(1.1)
CONTENT_LEFT = Inches(0.5)
CONTENT_WIDTH = Inches(12.3)
CONTENT_HEIGHT = Inches(5.7)


# ---------------------------------------------------------------------------
# Helper functions
# ---------------------------------------------------------------------------

def hex_to_rgb(hex_str: str) -> RGBColor:
    """
    Convert a hex color string to an RGBColor.

    Args:
        hex_str: Color string like "#FF6B35" or "FF6B35"

    Returns:
        RGBColor instance
    """
    hex_str = hex_str.lstrip("#")
    red = int(hex_str[0:2], 16)
    green = int(hex_str[2:4], 16)
    blue = int(hex_str[4:6], 16)
    return RGBColor(red, green, blue)


def interpolate_heatmap_color(
    value: float,
    min_val: float,
    max_val: float
) -> RGBColor:
    """
    Return a gradient color for heatmap cells based on value position
    between min_val and max_val.

    Uses the HEATMAP_STOPS gradient (white → peach → orange → red).
    Values outside [min_val, max_val] are clamped.

    Args:
        value: The data value to colorize
        min_val: Minimum value in the range (maps to white)
        max_val: Maximum value in the range (maps to red)

    Returns:
        RGBColor for the cell background
    """
    # Clamp and normalize to [0, 1]
    if max_val == min_val:
        normalized = 0.0
    else:
        normalized = max(0.0, min(1.0, (value - min_val) / (max_val - min_val)))

    # Find the two gradient stops that bracket this value
    for i in range(len(HEATMAP_STOPS) - 1):
        stop_low, color_low = HEATMAP_STOPS[i]
        stop_high, color_high = HEATMAP_STOPS[i + 1]

        if stop_low <= normalized <= stop_high:
            # Linear interpolation between the two stops
            segment_range = stop_high - stop_low
            if segment_range == 0:
                ratio = 0.0
            else:
                ratio = (normalized - stop_low) / segment_range

            red = int(color_low[0] + ratio * (color_high[0] - color_low[0]))
            green = int(color_low[1] + ratio * (color_high[1] - color_low[1]))
            blue = int(color_low[2] + ratio * (color_high[2] - color_low[2]))
            return RGBColor(red, green, blue)

    # Fallback: return the last stop color
    last_color = HEATMAP_STOPS[-1][1]
    return RGBColor(*last_color)


def add_title(slide, text: str) -> None:
    """
    Add a styled title text box at the top of a slide.

    Args:
        slide: pptx Slide object
        text: Title text to display
    """
    from pptx.util import Inches, Pt

    text_box = slide.shapes.add_textbox(
        TITLE_LEFT, TITLE_TOP, TITLE_WIDTH, TITLE_HEIGHT
    )
    text_frame = text_box.text_frame
    text_frame.word_wrap = True

    paragraph = text_frame.paragraphs[0]
    paragraph.text = text
    paragraph.font.name = FONT_TITLE
    paragraph.font.size = SIZE_TITLE
    paragraph.font.color.rgb = COLOR_TITLE
    paragraph.font.bold = True


def add_source_text(slide, text: str = "Source: Fruity Line store visit data") -> None:
    """
    Add 8pt source attribution text at the bottom-left of a slide.

    Args:
        slide: pptx Slide object
        text: Source text to display
    """
    text_box = slide.shapes.add_textbox(
        SOURCE_LEFT, SOURCE_TOP, SOURCE_WIDTH, SOURCE_HEIGHT
    )
    text_frame = text_box.text_frame

    paragraph = text_frame.paragraphs[0]
    paragraph.text = text
    paragraph.font.name = FONT_SOURCE
    paragraph.font.size = SIZE_SOURCE
    paragraph.font.color.rgb = COLOR_BODY_TEXT


def add_logo(slide, logo_path: Path | None = None) -> None:
    """
    Place the Fruity Line logo PNG at the bottom-right corner of a slide.

    Skips silently with a log warning if the logo file doesn't exist.

    Args:
        slide: pptx Slide object
        logo_path: Path to the logo PNG file. Defaults to DEFAULT_LOGO_PATH.
    """
    if logo_path is None:
        logo_path = DEFAULT_LOGO_PATH

    if not logo_path.exists():
        logger.warning(f"Logo file not found at {logo_path} — skipping logo")
        return

    slide.shapes.add_picture(
        str(logo_path), LOGO_LEFT, LOGO_TOP, LOGO_WIDTH, LOGO_HEIGHT
    )


def add_slide_number(slide, number: int) -> None:
    """
    Add a slide number at the bottom-right of a slide.

    Args:
        slide: pptx Slide object
        number: Slide number to display
    """
    text_box = slide.shapes.add_textbox(
        SLIDE_NUMBER_LEFT, SLIDE_NUMBER_TOP,
        SLIDE_NUMBER_WIDTH, SLIDE_NUMBER_HEIGHT
    )
    text_frame = text_box.text_frame

    paragraph = text_frame.paragraphs[0]
    paragraph.text = str(number)
    paragraph.font.name = FONT_BODY
    paragraph.font.size = SIZE_SLIDE_NUMBER
    paragraph.font.color.rgb = COLOR_BODY_TEXT
    paragraph.alignment = PP_ALIGN.RIGHT


def apply_chart_style(chart, colors: list[RGBColor] | None = None) -> None:
    """
    Apply consistent visual styling to a python-pptx chart object.

    Sets font, removes gridlines, and applies series colors.

    Args:
        chart: pptx Chart object
        colors: List of RGBColor for each series. Defaults to COLORS_CATEGORICAL.
    """
    if colors is None:
        colors = COLORS_CATEGORICAL

    # Style the chart font
    chart.font.name = FONT_BODY
    chart.font.size = SIZE_CHART_LABEL

    # Remove major gridlines on value axis for cleaner look
    value_axis = chart.value_axis
    value_axis.major_gridlines.format.line.fill.background()
    value_axis.has_minor_gridlines = False

    # Style axis labels
    value_axis.tick_labels.font.size = SIZE_CHART_LABEL
    value_axis.tick_labels.font.name = FONT_BODY

    category_axis = chart.category_axis
    category_axis.tick_labels.font.size = SIZE_CHART_LABEL
    category_axis.tick_labels.font.name = FONT_BODY

    # Apply series colors
    for idx, series in enumerate(chart.series):
        color = colors[idx % len(colors)]
        series.format.fill.solid()
        series.format.fill.fore_color.rgb = color


def style_table_header(table, header_color: RGBColor | None = None) -> None:
    """
    Style the first row of a pptx Table as a header.

    Args:
        table: pptx Table object
        header_color: Background color for header cells. Defaults to brand secondary.
    """
    if header_color is None:
        header_color = COLOR_BRAND_SECONDARY

    for cell in table.rows[0].cells:
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_color

        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.name = FONT_BODY
            paragraph.font.size = SIZE_TABLE_HEADER
            paragraph.font.bold = True
            paragraph.font.color.rgb = COLOR_WHITE


def style_table_body(table) -> None:
    """
    Apply consistent body styling to all non-header rows of a table.

    Args:
        table: pptx Table object
    """
    for row_idx in range(1, len(table.rows)):
        for cell in table.rows[row_idx].cells:
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.name = FONT_BODY
                paragraph.font.size = SIZE_TABLE_BODY
                paragraph.font.color.rgb = COLOR_BODY_TEXT


def add_subtitle(slide, text: str, left=None, top=None) -> None:
    """
    Add a smaller subtitle text below a chart or section.

    Args:
        slide: pptx Slide object
        text: Subtitle text
        left: Left position (defaults to CONTENT_LEFT)
        top: Top position (defaults to just below title)
    """
    if left is None:
        left = CONTENT_LEFT
    if top is None:
        top = Inches(0.85)

    text_box = slide.shapes.add_textbox(left, top, Inches(12.0), Inches(0.4))
    text_frame = text_box.text_frame

    paragraph = text_frame.paragraphs[0]
    paragraph.text = text
    paragraph.font.name = FONT_BODY
    paragraph.font.size = SIZE_SUBTITLE
    paragraph.font.color.rgb = COLOR_BODY_TEXT
