"""
Main PPTX generation orchestrator.

Takes the slide_data dict from analysis/slide_data.py generate_all_slide_data()
and produces a complete branded Fruity Line PowerPoint presentation.

Flow:
    slide_data dict → per-slide builder functions → styled .pptx file

Uses FL_template.pptx as the base when available, otherwise creates a
blank widescreen presentation. Each slide is built by a dedicated builder
function dispatched via the chart_type from config/storyline.py.
"""

import logging
from pathlib import Path
from typing import Any

import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt, Emu

from config.storyline import STORYLINE, get_slide_config
from output.style import (
    SLIDE_WIDTH,
    SLIDE_HEIGHT,
    CONTENT_TOP,
    CONTENT_LEFT,
    CONTENT_WIDTH,
    CONTENT_HEIGHT,
    COLORS_CATEGORICAL,
    COLORS_BINARY,
    DEFAULT_TEMPLATE_PATH,
    DEFAULT_LOGO_PATH,
    add_title,
    add_source_text,
    add_logo,
    add_slide_number,
)
from output.chart_builder import (
    add_grouped_bar_chart,
    add_bar_chart,
    add_heatmap_table,
    add_data_table,
)

logger = logging.getLogger(__name__)

# Source text that appears on every slide
SOURCE_TEXT = "Source: Fruity Line store visit data"


def generate_presentation(
    slide_data: dict[int, Any],
    template_path: Path | None = None,
    logo_path: Path | None = None,
    output_path: Path = Path("output/presentation.pptx"),
    headlines: dict[int, str] | None = None,
) -> Path:
    """
    Generate a complete branded PowerPoint presentation from slide data.

    Args:
        slide_data: Dict mapping slide_number (1-10) to analysis data,
            as returned by generate_all_slide_data()
        template_path: Path to FL_template.pptx. Defaults to assets/FL template.pptx.
        logo_path: Path to logo PNG. Defaults to assets/fruity_line_logo.png.
        output_path: Where to save the generated .pptx file.
        headlines: Dict mapping slide_number to headline string.
            If None, uses title_template from storyline config.

    Returns:
        Path to the generated .pptx file

    Raises:
        OSError: If the output directory cannot be created
    """
    if template_path is None:
        template_path = DEFAULT_TEMPLATE_PATH
    if logo_path is None:
        logo_path = DEFAULT_LOGO_PATH

    # Build fallback headlines from config if none provided
    if headlines is None:
        headlines = {
            sc["slide_number"]: sc["title_template"]
            for sc in STORYLINE
        }

    # Open template or create blank presentation
    presentation = _open_or_create_presentation(template_path)

    # Find the best layout to use for new slides
    slide_layout = _get_blank_layout(presentation)

    # Build each slide in storyline order
    for slide_config in STORYLINE:
        slide_num = slide_config["slide_number"]
        chart_type = slide_config["chart_type"]
        data = slide_data.get(slide_num)

        # Skip slides with no data (log warning)
        if data is None:
            logger.warning(
                f"Slide {slide_num} ({slide_config['title_template']}): "
                f"no data available, skipping"
            )
            continue

        headline = headlines.get(slide_num, slide_config["title_template"])

        # Add a new blank slide
        slide = presentation.slides.add_slide(slide_layout)

        # Dispatch to the appropriate builder
        try:
            if chart_type == "grouped_bar_grid":
                _build_slide_market_fingerprint(
                    slide, data, headline, slide_num, logo_path
                )
            elif chart_type == "heatmap":
                _build_slide_brand_landscape(
                    slide, data, headline, slide_num, logo_path
                )
            elif chart_type == "combo":
                _build_slide_retailer_sizing(
                    slide, data, headline, slide_num, logo_path
                )
            elif chart_type == "four_bar_grid":
                retailer = slide_config["params"].get("retailer", "")
                _build_slide_retailer_deep_dive(
                    slide, data, retailer, headline, slide_num, logo_path
                )
            else:
                logger.error(
                    f"Slide {slide_num}: unknown chart_type '{chart_type}'"
                )
        except Exception as exc:
            logger.error(
                f"Error building slide {slide_num}: {exc}",
                exc_info=True,
            )

    # Ensure output directory exists
    output_path.parent.mkdir(parents=True, exist_ok=True)

    presentation.save(str(output_path))
    logger.info(f"Presentation saved to {output_path} ({len(presentation.slides)} slides)")

    return output_path


# ---------------------------------------------------------------------------
# Template handling
# ---------------------------------------------------------------------------

def _open_or_create_presentation(template_path: Path) -> Presentation:
    """
    Open a template PPTX or create a blank widescreen presentation.

    If the template exists, opens it and removes all existing content slides
    (preserving slide masters and layouts). If not, creates a blank
    presentation with 13.33 x 7.5 inch slide dimensions.

    Args:
        template_path: Path to the template file

    Returns:
        A pptx Presentation object ready for new slides
    """
    if template_path.exists():
        logger.info(f"Loading template from {template_path}")
        presentation = Presentation(str(template_path))

        # Remove existing content slides (iterate in reverse to avoid index shift)
        slide_ids_to_remove = [slide.slide_id for slide in presentation.slides]
        for slide_id in slide_ids_to_remove:
            _remove_slide_by_id(presentation, slide_id)

        logger.info(
            f"Template loaded: {len(presentation.slide_layouts)} layouts available"
        )
    else:
        logger.warning(
            f"Template not found at {template_path} — creating blank presentation"
        )
        presentation = Presentation()
        presentation.slide_width = Emu(int(SLIDE_WIDTH))
        presentation.slide_height = Emu(int(SLIDE_HEIGHT))

    return presentation


def _remove_slide_by_id(presentation: Presentation, slide_id: int) -> None:
    """
    Remove a slide from a presentation by its slide ID.

    Uses low-level XML manipulation since python-pptx doesn't expose
    a public slide deletion API. Removes both the slide ID reference
    and the underlying relationship/part to avoid orphaned XML warnings.

    Args:
        presentation: The Presentation object
        slide_id: The numeric slide ID to remove
    """
    # Find and remove the sldId element from the slide list
    slide_list = presentation.slides._sldIdLst
    rel_id = None
    for sldId in slide_list:
        if int(sldId.get("id")) == slide_id:
            rel_id = sldId.get("{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id")
            slide_list.remove(sldId)
            break

    # Also remove the underlying relationship to avoid orphaned parts
    if rel_id is not None:
        try:
            presentation.part.drop_rel(rel_id)
        except Exception:
            pass


def _get_blank_layout(presentation: Presentation):
    """
    Find the best blank/empty layout from the presentation's slide masters.

    Looks for layouts named "Titel en leeg" (Dutch template) or "Blank"
    or similar. Falls back to the last layout if none match.

    Args:
        presentation: The Presentation object

    Returns:
        A SlideLayout object
    """
    preferred_names = [
        "titel en leeg",
        "title and blank",
        "blank",
        "leeg",
        "empty",
    ]

    for layout in presentation.slide_layouts:
        layout_name = layout.name.lower().strip()
        for preferred in preferred_names:
            if preferred in layout_name:
                logger.info(f"Using slide layout: '{layout.name}'")
                return layout

    # Fallback: use the last layout (often the most blank)
    fallback = presentation.slide_layouts[-1]
    logger.info(f"No preferred layout found, using: '{fallback.name}'")
    return fallback


# ---------------------------------------------------------------------------
# Per-slide builder functions
# ---------------------------------------------------------------------------

def _build_slide_market_fingerprint(
    slide,
    slide_data: dict[str, pd.DataFrame],
    headline: str,
    slide_num: int,
    logo_path: Path,
) -> None:
    """
    Build Slide 1: Market Fingerprint — 5 grouped bar charts side by side.

    Each chart shows a category breakdown (Product Type, PL vs Branded, etc.)
    across all retailers.

    Args:
        slide: pptx Slide object
        slide_data: Dict with keys like "product_type", "pl_vs_branded", etc.
        headline: Slide title text
        slide_num: Slide number for footer
        logo_path: Path to logo file
    """
    add_title(slide, headline)

    # Chart layout: 5 charts evenly spaced across the slide
    chart_width = Inches(2.3)
    chart_height = Inches(4.5)
    chart_top = CONTENT_TOP
    start_left = Inches(0.3)
    gap = Inches(0.25)

    # Map data keys to their display configs
    chart_configs = [
        ("product_type", "% by Product Type", "Product Type"),
        ("pl_vs_branded", "% PL vs Branded", "Branded/Private Label"),
        ("extraction", "% Cold Pressed vs Other", "Juice Extraction Method"),
        ("hpp", "% HPP vs Non-HPP", "HPP Treatment"),
        ("need_state", "% Indulgence vs Functional", "Need State"),
    ]

    for idx, (data_key, chart_title, category_col) in enumerate(chart_configs):
        chart_left = start_left + idx * (chart_width + gap)
        chart_df = slide_data.get(data_key, pd.DataFrame())

        # Pick color scheme: binary for 2-value categories, categorical otherwise
        unique_categories = (
            chart_df[category_col].nunique()
            if not chart_df.empty and category_col in chart_df.columns
            else 0
        )
        colors = COLORS_BINARY if unique_categories <= 2 else COLORS_CATEGORICAL

        add_grouped_bar_chart(
            slide=slide,
            data=chart_df,
            groupby_col="Retailer",
            category_col=category_col,
            value_col="Percentage",
            left=chart_left,
            top=chart_top,
            width=chart_width,
            height=chart_height,
            title=chart_title,
            colors=colors,
        )

    add_source_text(slide, SOURCE_TEXT)
    add_logo(slide, logo_path)
    add_slide_number(slide, slide_num)


def _build_slide_brand_landscape(
    slide,
    slide_data: pd.DataFrame,
    headline: str,
    slide_num: int,
    logo_path: Path,
) -> None:
    """
    Build Slide 2: Brand Landscape — heatmap table.

    Shows top brands as rows, retailers as columns, with conditional
    coloring on the retailer share values.

    Args:
        slide: pptx Slide object
        slide_data: DataFrame with Brand, Total Market Share, retailer
            columns, % Cold Pressed, % Functional
        headline: Slide title text
        slide_num: Slide number for footer
        logo_path: Path to logo file
    """
    add_title(slide, headline)

    # Retailer columns that get heatmap coloring
    heatmap_columns = [
        "Aldi", "Lidl", "M&S", "Sainsbury's",
        "Tesco", "Tesco Express", "Waitrose",
    ]

    # Table spans most of the content area
    table_left = Inches(0.3)
    table_top = CONTENT_TOP
    table_width = Inches(12.7)
    table_height = Inches(5.5)

    add_heatmap_table(
        slide=slide,
        data=slide_data,
        left=table_left,
        top=table_top,
        width=table_width,
        height=table_height,
        heatmap_columns=heatmap_columns,
    )

    add_source_text(slide, SOURCE_TEXT)
    add_logo(slide, logo_path)
    add_slide_number(slide, slide_num)


def _build_slide_retailer_sizing(
    slide,
    slide_data: tuple[pd.DataFrame, pd.DataFrame],
    headline: str,
    slide_num: int,
    logo_path: Path,
) -> None:
    """
    Build Slide 3: Retailer Sizing — table + grouped bar chart.

    Left half: table with Retailer, Store Format, Avg SKU Count, Avg Facings.
    Right half: grouped bar chart showing % PL, % HPP, % Cold Pressed per retailer.

    Args:
        slide: pptx Slide object
        slide_data: Tuple of (table_data, chart_data) DataFrames
        headline: Slide title text
        slide_num: Slide number for footer
        logo_path: Path to logo file
    """
    add_title(slide, headline)

    table_data, chart_data = slide_data

    # Left half: data table
    table_left = Inches(0.3)
    table_top = CONTENT_TOP
    table_width = Inches(5.5)
    table_height = Inches(5.0)

    add_data_table(
        slide=slide,
        data=table_data,
        left=table_left,
        top=table_top,
        width=table_width,
        height=table_height,
    )

    # Right half: grouped bar chart with 3 metrics as series
    # Reshape chart_data for grouped bar: melt % PL, % HPP, % Cold Pressed
    if not chart_data.empty:
        chart_melted = chart_data.melt(
            id_vars=["Retailer"],
            value_vars=["% PL", "% HPP", "% Cold Pressed"],
            var_name="Metric",
            value_name="Percentage",
        )

        chart_left = Inches(6.2)
        chart_top = CONTENT_TOP
        chart_width = Inches(6.5)
        chart_height = Inches(5.0)

        add_grouped_bar_chart(
            slide=slide,
            data=chart_melted,
            groupby_col="Retailer",
            category_col="Metric",
            value_col="Percentage",
            left=chart_left,
            top=chart_top,
            width=chart_width,
            height=chart_height,
            title="Category Shares by Retailer",
            colors=COLORS_CATEGORICAL[:3],
        )

    add_source_text(slide, SOURCE_TEXT)
    add_logo(slide, logo_path)
    add_slide_number(slide, slide_num)


def _build_slide_retailer_deep_dive(
    slide,
    slide_data: dict[str, pd.DataFrame],
    retailer: str,
    headline: str,
    slide_num: int,
    logo_path: Path,
) -> None:
    """
    Build Slides 4-10: Retailer Deep Dive — 4 bar charts in a 2x2 grid.

    Top-left: Product Type (absolute facings)
    Top-right: PL vs Branded (%)
    Bottom-left: Extraction Method (%)
    Bottom-right: Need State (%)

    Args:
        slide: pptx Slide object
        slide_data: Dict with keys "product_type", "pl_vs_branded",
            "extraction", "need_state"
        retailer: Retailer name for context
        headline: Slide title text
        slide_num: Slide number for footer
        logo_path: Path to logo file
    """
    add_title(slide, headline)

    # 2x2 grid positions
    chart_width = Inches(5.8)
    chart_height = Inches(2.7)
    left_col = Inches(0.4)
    right_col = Inches(6.7)
    top_row = CONTENT_TOP
    bottom_row = Inches(3.9)

    # Chart configs: (data_key, title, category_col, value_col, number_format, use_stacked)
    chart_configs = [
        (
            "product_type", "Facings by Product Type",
            "Product Type", "Facings", "0", False,
        ),
        (
            "pl_vs_branded", "% PL vs Branded",
            "Branded/Private Label", "Percentage", "0.0", True,
        ),
        (
            "extraction", "% Cold Pressed vs Other",
            "Juice Extraction Method", "Percentage", "0.0", True,
        ),
        (
            "need_state", "% Indulgence vs Functional",
            "Need State", "Percentage", "0.0", True,
        ),
    ]

    positions = [
        (left_col, top_row),
        (right_col, top_row),
        (left_col, bottom_row),
        (right_col, bottom_row),
    ]

    for (data_key, chart_title, cat_col, val_col, num_fmt, stacked), (pos_left, pos_top) in zip(
        chart_configs, positions
    ):
        chart_df = slide_data.get(data_key, pd.DataFrame())

        # Pick colors: binary for 2-value categories, categorical otherwise
        num_categories = len(chart_df) if not chart_df.empty else 0
        colors = COLORS_BINARY if num_categories <= 2 else COLORS_CATEGORICAL

        add_bar_chart(
            slide=slide,
            data=chart_df,
            category_col=cat_col,
            value_col=val_col,
            left=pos_left,
            top=pos_top,
            width=chart_width,
            height=chart_height,
            title=chart_title,
            colors=colors,
            show_data_labels=True,
            number_format=num_fmt,
            use_stacked=stacked,
        )

    add_source_text(slide, SOURCE_TEXT)
    add_logo(slide, logo_path)
    add_slide_number(slide, slide_num)
