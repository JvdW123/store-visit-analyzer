"""
Tests for the Phase 2 PowerPoint generation pipeline.

Tests cover:
- Valid .pptx file generation
- Correct slide count
- Title presence on every slide
- Chart/table presence on expected slides
- Graceful handling of missing data
- Fallback when no template file
- Headline fallback without API key
"""

import os
from pathlib import Path

import pandas as pd
import pytest
from pptx import Presentation
from pptx.chart.chart import Chart
from pptx.table import Table

from output.pptx_generator import generate_presentation
from output.headline_generator import generate_all_headlines
from output.style import (
    hex_to_rgb,
    interpolate_heatmap_color,
    DEFAULT_TEMPLATE_PATH,
)
from config.storyline import STORYLINE


# ---------------------------------------------------------------------------
# Test fixtures — mock data matching calculations.py return schemas
# ---------------------------------------------------------------------------

RETAILERS = ["Aldi", "Lidl", "M&S", "Sainsbury's", "Tesco", "Tesco Express", "Waitrose"]


def _make_share_by_category_df(
    groupby_col: str,
    category_col: str,
    categories: list[str],
) -> pd.DataFrame:
    """Create a mock DataFrame matching share_by_category() output."""
    rows = []
    for retailer in RETAILERS:
        total = 100.0
        per_cat = total / len(categories)
        for cat in categories:
            rows.append({
                groupby_col: retailer,
                category_col: cat,
                "Count": int(per_cat),
                "Percentage": per_cat,
            })
    return pd.DataFrame(rows)


def _make_market_fingerprint_data() -> dict[str, pd.DataFrame]:
    """Mock data for slide 1 (market_fingerprint return type)."""
    return {
        "product_type": _make_share_by_category_df(
            "Retailer", "Product Type",
            ["Pure Juices", "Smoothies", "Shots"],
        ),
        "pl_vs_branded": _make_share_by_category_df(
            "Retailer", "Branded/Private Label",
            ["Private Label", "Branded"],
        ),
        "extraction": _make_share_by_category_df(
            "Retailer", "Juice Extraction Method",
            ["Cold Pressed", "Centrifugal"],
        ),
        "hpp": _make_share_by_category_df(
            "Retailer", "HPP Treatment",
            ["Yes", "No"],
        ),
        "need_state": _make_share_by_category_df(
            "Retailer", "Need State",
            ["Indulgence", "Functional"],
        ),
    }


def _make_brand_heatmap_data() -> pd.DataFrame:
    """Mock data for slide 2 (brand_retailer_heatmap return type)."""
    brands = [f"Brand_{i}" for i in range(1, 11)]
    rows = []
    for brand in brands:
        row = {
            "Brand": brand,
            "Total Market Share": round(10.0 / len(brands), 1),
        }
        for retailer in RETAILERS:
            row[retailer] = round(100.0 / len(brands), 1)
        row["% Cold Pressed"] = 25.0
        row["% Functional"] = 15.0
        rows.append(row)
    return pd.DataFrame(rows)


def _make_retailer_sizing_data() -> tuple[pd.DataFrame, pd.DataFrame]:
    """Mock data for slide 3 (retailer_sizing return type)."""
    table_data = pd.DataFrame({
        "Retailer": RETAILERS,
        "Store Format": ["Standard"] * len(RETAILERS),
        "Avg SKU Count": [30] * len(RETAILERS),
        "Avg Facings": [45] * len(RETAILERS),
    })
    chart_data = pd.DataFrame({
        "Retailer": RETAILERS,
        "% PL": [40.0] * len(RETAILERS),
        "% HPP": [15.0] * len(RETAILERS),
        "% Cold Pressed": [20.0] * len(RETAILERS),
    })
    return table_data, chart_data


def _make_retailer_deep_dive_data() -> dict[str, pd.DataFrame]:
    """Mock data for slides 4-10 (retailer_deep_dive return type)."""
    return {
        "product_type": pd.DataFrame({
            "Product Type": ["Pure Juices", "Smoothies", "Shots"],
            "Facings": [40, 30, 20],
        }),
        "pl_vs_branded": pd.DataFrame({
            "Branded/Private Label": ["Private Label", "Branded"],
            "Percentage": [55.0, 45.0],
        }),
        "extraction": pd.DataFrame({
            "Juice Extraction Method": ["Cold Pressed", "Centrifugal"],
            "Percentage": [60.0, 40.0],
        }),
        "need_state": pd.DataFrame({
            "Need State": ["Indulgence", "Functional"],
            "Percentage": [65.0, 35.0],
        }),
    }


def _make_complete_slide_data() -> dict[int, object]:
    """Build a complete slide_data dict for all 10 slides."""
    slide_data: dict[int, object] = {
        1: _make_market_fingerprint_data(),
        2: _make_brand_heatmap_data(),
        3: _make_retailer_sizing_data(),
    }
    for slide_num in range(4, 11):
        slide_data[slide_num] = _make_retailer_deep_dive_data()
    return slide_data


# ---------------------------------------------------------------------------
# Tests for pptx_generator.py
# ---------------------------------------------------------------------------

class TestGeneratePresentation:
    """Tests for the main generate_presentation function."""

    def test_generates_valid_pptx_file(self, tmp_path: Path) -> None:
        """The output file should exist and be openable by python-pptx."""
        slide_data = _make_complete_slide_data()
        output_file = tmp_path / "test_output.pptx"

        result_path = generate_presentation(
            slide_data=slide_data,
            template_path=Path("nonexistent_template.pptx"),
            logo_path=Path("nonexistent_logo.png"),
            output_path=output_file,
        )

        assert result_path.exists()
        # Verify it's a valid pptx by opening it
        presentation = Presentation(str(result_path))
        assert presentation is not None

    def test_slide_count_matches_storyline(self, tmp_path: Path) -> None:
        """Output should have exactly 10 slides (one per storyline entry)."""
        slide_data = _make_complete_slide_data()
        output_file = tmp_path / "test_count.pptx"

        result_path = generate_presentation(
            slide_data=slide_data,
            template_path=Path("nonexistent_template.pptx"),
            logo_path=Path("nonexistent_logo.png"),
            output_path=output_file,
        )

        presentation = Presentation(str(result_path))
        assert len(presentation.slides) == len(STORYLINE)

    def test_each_slide_has_title(self, tmp_path: Path) -> None:
        """Every slide should contain at least one text shape with content."""
        slide_data = _make_complete_slide_data()
        output_file = tmp_path / "test_titles.pptx"

        result_path = generate_presentation(
            slide_data=slide_data,
            template_path=Path("nonexistent_template.pptx"),
            logo_path=Path("nonexistent_logo.png"),
            output_path=output_file,
        )

        presentation = Presentation(str(result_path))
        for slide_idx, slide in enumerate(presentation.slides):
            has_text = any(
                shape.has_text_frame and shape.text_frame.text.strip()
                for shape in slide.shapes
            )
            assert has_text, f"Slide {slide_idx + 1} has no text shapes"

    def test_heatmap_slide_has_table(self, tmp_path: Path) -> None:
        """Slide 2 (Brand Landscape) should contain a Table shape."""
        slide_data = _make_complete_slide_data()
        output_file = tmp_path / "test_heatmap.pptx"

        result_path = generate_presentation(
            slide_data=slide_data,
            template_path=Path("nonexistent_template.pptx"),
            logo_path=Path("nonexistent_logo.png"),
            output_path=output_file,
        )

        presentation = Presentation(str(result_path))
        # Slide 2 is at index 1 (0-based)
        slide_2 = presentation.slides[1]
        has_table = any(shape.has_table for shape in slide_2.shapes)
        assert has_table, "Slide 2 should contain a table (heatmap)"

    def test_deep_dive_slides_have_charts(self, tmp_path: Path) -> None:
        """Slides 4-10 should each contain chart shapes."""
        slide_data = _make_complete_slide_data()
        output_file = tmp_path / "test_charts.pptx"

        result_path = generate_presentation(
            slide_data=slide_data,
            template_path=Path("nonexistent_template.pptx"),
            logo_path=Path("nonexistent_logo.png"),
            output_path=output_file,
        )

        presentation = Presentation(str(result_path))

        # Slides 4-10 are at indices 3-9 (0-based)
        for slide_idx in range(3, 10):
            slide = presentation.slides[slide_idx]
            has_chart = any(shape.has_chart for shape in slide.shapes)
            assert has_chart, (
                f"Slide {slide_idx + 1} (retailer deep dive) should contain charts"
            )

    def test_missing_slide_data_skips_gracefully(self, tmp_path: Path) -> None:
        """If a slide's data is None, it should be skipped without crashing."""
        slide_data = _make_complete_slide_data()
        # Remove data for slide 3
        slide_data[3] = None
        output_file = tmp_path / "test_missing.pptx"

        result_path = generate_presentation(
            slide_data=slide_data,
            template_path=Path("nonexistent_template.pptx"),
            logo_path=Path("nonexistent_logo.png"),
            output_path=output_file,
        )

        presentation = Presentation(str(result_path))
        # Should have 9 slides (one skipped)
        assert len(presentation.slides) == 9

    def test_no_template_creates_blank_presentation(self, tmp_path: Path) -> None:
        """Without a template file, a valid blank presentation is created."""
        slide_data = _make_complete_slide_data()
        output_file = tmp_path / "test_blank.pptx"

        result_path = generate_presentation(
            slide_data=slide_data,
            template_path=Path("definitely_does_not_exist.pptx"),
            logo_path=Path("nonexistent_logo.png"),
            output_path=output_file,
        )

        assert result_path.exists()
        presentation = Presentation(str(result_path))
        assert len(presentation.slides) == 10

    def test_custom_headlines_appear_on_slides(self, tmp_path: Path) -> None:
        """Custom headlines should appear in the slide text."""
        slide_data = _make_complete_slide_data()
        custom_headlines = {
            i: f"Custom Headline {i}" for i in range(1, 11)
        }
        output_file = tmp_path / "test_headlines.pptx"

        result_path = generate_presentation(
            slide_data=slide_data,
            template_path=Path("nonexistent_template.pptx"),
            logo_path=Path("nonexistent_logo.png"),
            output_path=output_file,
            headlines=custom_headlines,
        )

        presentation = Presentation(str(result_path))
        for slide_idx, slide in enumerate(presentation.slides):
            slide_num = slide_idx + 1
            expected_headline = f"Custom Headline {slide_num}"

            all_text = " ".join(
                shape.text_frame.text
                for shape in slide.shapes
                if shape.has_text_frame
            )
            assert expected_headline in all_text, (
                f"Slide {slide_num} should contain headline '{expected_headline}'"
            )

    def test_with_real_template_if_available(self, tmp_path: Path) -> None:
        """If the real FL template exists, test that it works."""
        if not DEFAULT_TEMPLATE_PATH.exists():
            pytest.skip("FL template not available")

        slide_data = _make_complete_slide_data()
        output_file = tmp_path / "test_real_template.pptx"

        result_path = generate_presentation(
            slide_data=slide_data,
            template_path=DEFAULT_TEMPLATE_PATH,
            logo_path=Path("nonexistent_logo.png"),
            output_path=output_file,
        )

        presentation = Presentation(str(result_path))
        assert len(presentation.slides) == 10


# ---------------------------------------------------------------------------
# Tests for headline_generator.py
# ---------------------------------------------------------------------------

class TestHeadlineGenerator:
    """Tests for headline generation with fallback behavior."""

    def test_fallback_without_api_key(self) -> None:
        """Without an API key, should return title_template for every slide."""
        slide_data = _make_complete_slide_data()
        headlines = generate_all_headlines(slide_data, api_key=None)

        assert len(headlines) == 10
        for slide_config in STORYLINE:
            slide_num = slide_config["slide_number"]
            assert headlines[slide_num] == slide_config["title_template"]

    def test_fallback_with_empty_api_key(self) -> None:
        """An empty string API key should also trigger fallback."""
        slide_data = _make_complete_slide_data()
        headlines = generate_all_headlines(slide_data, api_key="")

        assert len(headlines) == 10
        for slide_config in STORYLINE:
            slide_num = slide_config["slide_number"]
            assert headlines[slide_num] == slide_config["title_template"]


# ---------------------------------------------------------------------------
# Tests for style.py utilities
# ---------------------------------------------------------------------------

class TestStyleHelpers:
    """Tests for style utility functions."""

    def test_hex_to_rgb_with_hash(self) -> None:
        """Should parse hex string with leading #."""
        color = hex_to_rgb("#FF6B35")
        assert color == (0xFF, 0x6B, 0x35)

    def test_hex_to_rgb_without_hash(self) -> None:
        """Should parse hex string without leading #."""
        color = hex_to_rgb("004E89")
        assert color == (0x00, 0x4E, 0x89)

    def test_heatmap_color_at_minimum(self) -> None:
        """Minimum value should return white (or near-white)."""
        color = interpolate_heatmap_color(0.0, 0.0, 100.0)
        # Should be white: RGB(255, 255, 255)
        assert color == (255, 255, 255)

    def test_heatmap_color_at_maximum(self) -> None:
        """Maximum value should return the darkest red."""
        color = interpolate_heatmap_color(100.0, 0.0, 100.0)
        # Should be dark red: RGB(199, 62, 29)
        assert color == (199, 62, 29)

    def test_heatmap_color_equal_min_max(self) -> None:
        """When min == max, should return white (normalized = 0)."""
        color = interpolate_heatmap_color(50.0, 50.0, 50.0)
        assert color == (255, 255, 255)

    def test_heatmap_color_clamping(self) -> None:
        """Values outside range should be clamped, not crash."""
        # Value below minimum
        color_low = interpolate_heatmap_color(-10.0, 0.0, 100.0)
        assert color_low == (255, 255, 255)

        # Value above maximum
        color_high = interpolate_heatmap_color(200.0, 0.0, 100.0)
        assert color_high == (199, 62, 29)
